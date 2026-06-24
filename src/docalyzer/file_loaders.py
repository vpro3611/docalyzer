from __future__ import annotations

import csv
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Callable, Dict

try:
    import docx  # type: ignore
except ImportError:
    docx = None

try:
    import yaml  # type: ignore
except ImportError:
    yaml = None

try:
    from openpyxl import load_workbook  # type: ignore
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    Presentation = None

try:
    import PyPDF2  # type: ignore
except ImportError:
    PyPDF2 = None

try:
    from bs4 import BeautifulSoup  # type: ignore
except ImportError:
    BeautifulSoup = None

_LOADER_REGISTRY: Dict[str, Callable[[Path], str]] = {}


class FileLoadError(Exception):
    pass


class UnsupportedFileTypeError(FileLoadError):
    pass


class MissingDependencyError(FileLoadError):
    pass


def register_loader(extensions: tuple[str, ...]):
    def decorator(func: Callable[[Path], str]) -> Callable[[Path], str]:
        for extension in extensions:
            _LOADER_REGISTRY[extension] = func
        return func

    return decorator


def load_file(path: Path) -> str:
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    loader = _LOADER_REGISTRY.get(path.suffix.lower())
    if loader is None:
        raise UnsupportedFileTypeError(
            f"Unsupported file type: {path.suffix}. "
            f"Supported: {', '.join(sorted(_LOADER_REGISTRY))}"
        )

    return loader(path)


@register_loader((".txt", ".md", ".py", ".js", ".java", ".cpp"))
def load_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


@register_loader((".html",))
def load_html_file(path: Path) -> str:
    raw_text = path.read_text(encoding="utf-8", errors="replace")
    if BeautifulSoup is not None:
        soup = BeautifulSoup(raw_text, "html.parser")
        return soup.get_text(separator="\n", strip=True)

    text = re.sub(r"<script[\s\S]*?</script>", "", raw_text, flags=re.IGNORECASE)
    text = re.sub(r"<style[\s\S]*?</style>", "", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    return text


@register_loader((".json",))
def load_json_file(path: Path) -> str:
    raw_text = path.read_text(encoding="utf-8", errors="replace")
    parsed = json.loads(raw_text)
    return json.dumps(parsed, indent=2, ensure_ascii=False)


@register_loader((".xml",))
def load_xml_file(path: Path) -> str:
    raw_text = path.read_text(encoding="utf-8", errors="replace")
    try:
        root = ET.fromstring(raw_text)
    except ET.ParseError:
        return raw_text

    def _pretty(element: ET.Element, level: int = 0) -> str:
        indent = "  " * level
        attrs = "".join(
            f' {k}="{v}"' for k, v in element.attrib.items()
        )
        lines = [f"{indent}<{element.tag}{attrs}>"]
        if element.text and element.text.strip():
            lines.append(f"{indent}  {element.text.strip()}")
        for child in element:
            lines.append(_pretty(child, level + 1))
        lines.append(f"{indent}</{element.tag}>")
        return "\n".join(lines)

    return _pretty(root)


@register_loader((".yaml", ".yml"))
def load_yaml_file(path: Path) -> str:
    if yaml is None:
        raise MissingDependencyError(
            "Install PyYAML to read .yaml/.yml files: pip install pyyaml"
        )
    raw_text = path.read_text(encoding="utf-8", errors="replace")
    parsed = yaml.safe_load(raw_text)
    return yaml.safe_dump(parsed, sort_keys=False)


@register_loader((".csv",))
def load_csv_file(path: Path) -> str:
    with path.open(newline="", encoding="utf-8", errors="replace") as csvfile:
        reader = csv.reader(csvfile)
        lines = []
        for index, row in enumerate(reader):
            lines.append(", ".join(cell for cell in row if cell != ""))
            if index >= 19:
                break
    return "\n".join(lines)


@register_loader((".ipynb",))
def load_ipynb_file(path: Path) -> str:
    raw_json = path.read_text(encoding="utf-8", errors="replace")
    notebook = json.loads(raw_json)
    cells = notebook.get("cells", [])
    contents = []
    for cell in cells:
        cell_type = cell.get("cell_type")
        source = cell.get("source", [])
        source_text = "".join(source)
        if not source_text.strip():
            continue
        if cell_type == "code":
            contents.append("# code cell")
            contents.append(source_text)
        else:
            contents.append(source_text)
    return "\n\n".join(contents)


def _load_docx(path: Path) -> str:
    if docx is None:
        raise MissingDependencyError(
            "Install python-docx to read .docx files: pip install python-docx"
        )
    document = docx.Document(path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n\n".join(paragraphs)


@register_loader((".docx",))
def load_docx_file(path: Path) -> str:
    return _load_docx(path)


def _load_xlsx(path: Path) -> str:
    if load_workbook is None:
        raise MissingDependencyError(
            "Install openpyxl to read .xlsx files: pip install openpyxl"
        )
    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        rows = []
        for row_index, row in enumerate(
            sheet.iter_rows(values_only=True), start=1
        ):
            rows.append(", ".join(str(cell) for cell in row if cell is not None))
            if row_index >= 20:
                break
        if rows:
            parts.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
    return "\n\n".join(parts)


@register_loader((".xlsx",))
def load_xlsx_file(path: Path) -> str:
    return _load_xlsx(path)


def _load_pptx(path: Path) -> str:
    if Presentation is None:
        raise MissingDependencyError(
            "Install python-pptx to read .pptx files: pip install python-pptx"
        )
    presentation = Presentation(path)
    slides = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        if texts:
            slides.append(f"Slide {slide_index}: " + "\n".join(texts))
    return "\n\n".join(slides)


@register_loader((".pptx",))
def load_pptx_file(path: Path) -> str:
    return _load_pptx(path)


def _load_pdf(path: Path) -> str:
    if PyPDF2 is None:
        raise MissingDependencyError(
            "Install PyPDF2 to read .pdf files: pip install PyPDF2"
        )
    reader = PyPDF2.PdfReader(path)
    texts = []
    for page in reader.pages:
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            texts.append(page_text)
    return "\n\n".join(texts)


@register_loader((".pdf",))
def load_pdf_file(path: Path) -> str:
    return _load_pdf(path)


SUPPORTED_EXTENSIONS = tuple(sorted(_LOADER_REGISTRY.keys()))
